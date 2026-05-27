"""
生成 Assignment 2 演示 PPT (15 页)。

运行: PYTHONPATH=. .venv/bin/python demo/generate_ppt.py
输出: demo/presentation.pptx
"""

from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 主题色
DARK = RGBColor(0x1D, 0x1D, 0x1F)
BLUE = RGBColor(0x00, 0x71, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x86, 0x86, 0x8B)
GREEN = RGBColor(0x34, 0xC7, 0x59)
ORANGE = RGBColor(0xFF, 0x95, 0x00)
PURPLE = RGBColor(0xAF, 0x52, 0xDE)
RED = RGBColor(0xFF, 0x2D, 0x55)


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER


def _add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT

    # 蓝线
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), Inches(1.3), Inches(2), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # 要点
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)
        p.level = 0


def _add_arch_slide(prs: Presentation, title: str, content: str) -> None:
    """代码/架构展示页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(13)
    p2.font.color.rgb = WHITE
    p2.font.name = "Courier New"


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: 封面 ──
    _add_title_slide(prs, "AutoTestDesign", "AI 驱动的测试用例自动生成工具\nAssignment 2 成果演示")

    # ── Slide 2: 项目背景与目标 ──
    _add_content_slide(prs, "项目背景与目标", [
        "背景: 软件测试用例编写耗时且容易遗漏覆盖维度",
        "目标: 构建 AI 驱动的 AutoTestDesign 工具",
        "  输入: 结构化/自然语言需求 (CSV/文本/手动)",
        "  处理: FR1-FR7 全链路流水线",
        "  输出: 覆盖 EP/BVA/DT/STT 四种技术的测试套件",
        "基于 ISTQB Foundation Level + ISO/IEC/IEEE 29119-4 标准",
        "目标应用: AuthSystem — 用户注册与登录系统",
    ])

    # ── Slide 3: 工具架构总览 ──
    _add_arch_slide(prs, "AutoTestDesign 架构总览",
        "autotestdesign/\n"
        "├── app.py              # Streamlit 入口 (4 Tab)\n"
        "├── config.py           # 配置加载 (.env / config.json)\n"
        "├── core/               # 纯算法层 (不依赖 UI 或 LLM)\n"
        "│   ├── models.py       # Pydantic 数据模型\n"
        "│   ├── parsing.py      # FR1 需求解析\n"
        "│   ├── risk.py         # FR2 风险打分\n"
        "│   ├── blackbox/       # FR3 EP / BVA / DT\n"
        "│   ├── whitebox/       # FR4 状态机覆盖\n"
        "│   ├── oracle.py       # FR5 测试预言\n"
        "│   ├── export.py       # FR6 导出\n"
        "│   ├── optimizer.py    # FR7 套件优化\n"
        "│   └── pipeline.py     # 流水线编排\n"
        "├── llm/                # OpenAI 兼容客户端 + Prompt 模板\n"
        "├── ui/                 # Streamlit UI 组件\n"
        "└── examples/           # 内置示例需求"
    )

    # ── Slide 4: FR1-FR7 功能展示 ──
    _add_content_slide(prs, "FR1-FR7 功能需求覆盖", [
        "FR1: 需求输入与解析 — CSV (pandas) + 纯文本 + LLM 结构化",
        "FR2: 风险分析与优先级 — 中英双语关键词规则 + LLM rationale",
        "FR3: 黑盒测试设计 — EP/ BVA / DT 三种 ISO 29119-4 技术",
        "FR4: 白盒状态转换 — DSL 解析 + All-States + All-Transitions 覆盖",
        "FR5: 测试预言合成 — 越界算法兜底 + LLM 合成 + 失败降级",
        "FR6: 导出 — JSON / CSV / Excel (openpyxl)",
        "FR7: 测试套件优化 — 指纹去重 + 加权集合覆盖 + 优先级排序",
    ])

    # ── Slide 5: 目标应用 AuthSystem ──
    _add_content_slide(prs, "目标应用: AuthSystem", [
        "纯 Python 内存实现的用户注册与登录系统",
        "核心功能: register / login / logout / reset_password",
        "6 条业务规则:",
        "  BR-01: 年龄 18-120    BR-02: 密码 8-32 + 数字 + 字母",
        "  BR-03: 邮箱含 @ 和 .   BR-04: 用户名 3-20 且唯一",
        "  BR-05: 5 次失败 → 锁定  BR-06: 密码重置可解锁",
        "状态机: Guest → Registered → LoggedIn → Locked → Registered",
        "零外部依赖，dict 存储，~140 行代码",
    ])

    # ── Slide 6: 用工具测试 AuthSystem ──
    _add_content_slide(prs, "用 AutoTestDesign 测试 AuthSystem 的完整流程", [
        "① 准备 requirements.csv (6 条中英双语需求)",
        "② 加载需求到 AutoTestDesign (上传或选择文件)",
        "③ 配置覆盖策略 (EP/BVA/DT/STT 全选)",
        "④ 定义状态机 DSL (4 节点 8 转换)",
        "⑤ 运行 FR1→FR7 完整流水线 (关闭 LLM，规则兜底)",
        "⑥ 审查覆盖项 — 确认 6 条需求在各技术下的覆盖率",
        "⑦ 交互式编辑 — 调整优先级、添加覆盖项",
        "⑧ 导出 JSON/CSV/Excel 到 target_app/output/",
        "生成结果: 82 条用例 (EP 18 / BVA 30 / DT 24 / STT 10)",
    ])

    # ── Slide 7: 交互式设计审查 ──
    _add_content_slide(prs, "交互式设计审查 (核心增量)", [
        "阶段 2 新增的编辑能力:",
        "  覆盖项概览面板: 按需求×技术显示用例分布矩阵",
        "  策略编辑: 运行前开关各项测试技术",
        "  可编辑用例网格 (st.data_editor):",
        "    - 修改输入值、预期结果、优先级、标签",
        "    - 动态添加/删除行",
        "    - 同步修改到测试套件",
        "    - 一键导出含编辑的 CSV",
        "  覆盖项手动添加: 补充缺失的测试维度",
        "体现 Human-in-the-loop 测试设计最佳实践",
    ])

    # ── Slide 8: 风险分析报告要点 ──
    _add_content_slide(prs, "风险分析报告要点", [
        "识别 14 项风险: 3 高 / 8 中 / 3 低",
        "高优先级:",
        "  R-01 暴力破解攻击 (L4×I5=20)",
        "  R-05 弱密码策略绕过 (L4×I4=16)",
        "  R-09 密码明文内存存储 (L4×I5=20)",
        "缓解后: 高风险 3→0, 中风险 8→4",
        "核心建议: bcrypt 密码哈希, 统一错误消息, 速率限制",
    ])

    # ── Slide 9: 测试计划要点 ──
    _add_content_slide(prs, "测试计划要点", [
        "5 个测试套件: TS-1~TS-4 由工具生成, TS-5 为 pytest 手工用例",
        "测试范围: AuthSystem 全部公开 API + 校验层 + 状态机",
        "覆盖矩阵: 6 条需求 × 4 种技术 × 可追溯性",
        "测试框架: pytest 9.0 + parametrize + fixture",
        "成本估算: 使用 AutoTestDesign 节省 64% 工时",
        "  (手工 12 人日 → 工具辅助 4.3 人日)",
        "组织架构: 5 角色 × 7 职责覆盖开发/测试/文档/演示",
    ])

    # ── Slide 10: 详细测试设计要点 ──
    _add_content_slide(prs, "详细测试设计要点", [
        "Mainly 完整流程: 概念→覆盖项→策略→用例→Prompt→审查→结果→改进",
        "覆盖项识别: 6 需求 × 4-5 字段 = 30 个覆盖维度",
        "EP: 18 条, 按字段类型划分等价类 (int/string)",
        "BVA: 30 条, 5 个有界字段 × 6 个边界值",
        "DT: 24 条, 2^1 + 2^3 + 2^2 + 2^2 + 2^1 + 2^2 组合",
        "STT: 10 条, All-States (4 节点) + All-Transitions (8 转换)",
        "可追溯性: 每条用例关联需求 ID + 技术 + 字段",
        "交互式审查: 策略调整 + 覆盖项补充 + 优先级修正",
    ])

    # ── Slide 11: 测试执行结果 ──
    _add_content_slide(prs, "测试执行结果", [
        "pytest 执行: 49 个测试用例全部通过 (0.03s)",
        "AutoTestDesign 生成: 82 条用例 (0.5s 流水线)",
        "覆盖统计:",
        "  BR-01 年龄: 8 tests PASSED",
        "  BR-02 密码: 10 tests PASSED",
        "  BR-03 邮箱: 9 tests PASSED",
        "  BR-04 用户名: 6 tests PASSED",
        "  BR-05 锁定: 4 tests PASSED",
        "  BR-06 重置: 1 test PASSED",
        "  状态机转换: 6 tests PASSED",
        "  边缘场景: 5 tests PASSED",
        "AutoTestDesign 工具测试: 56 个单元/集成测试全部通过",
        "总计: 105 tests ALL PASSED",
    ])

    # ── Slide 12: 成本估算对比 ──
    _add_content_slide(prs, "成本估算对比: 手工 vs AutoTestDesign", [
        "| 活动                     | 手工(人日) | 工具(人日) | 节省   |",
        "|--------------------------|-----------|-----------|--------|",
        "| 需求分析与策略设计        | 1.0       | 0.5       | 50%    |",
        "| EP 用例编写 (18条)        | 0.5       | 0.1       | 80%    |",
        "| BVA 用例编写 (30条)       | 1.0       | 0.1       | 90%    |",
        "| DT 用例编写 (24条)        | 1.0       | 0.1       | 90%    |",
        "| STT 用例编写 (10条)       | 1.5       | 0.2       | 87%    |",
        "| pytest 脚本编写           | 1.0       | 0.5       | 50%    |",
        "| 测试执行与分析            | 0.5       | 0.3       | 40%    |",
        "| 文档编写 (3份报告)        | 4.5       | 2.3       | 49%    |",
        "| 总计                     | 12.0      | 4.3       | 64%    |",
        "",
        "基础设施成本: USD 0.50 (LLM API tokens)",
        "工具链: 全部 MIT/BSD 许可, 学术使用免费",
    ])

    # ── Slide 13: 总结与收获 ──
    _add_content_slide(prs, "总结与收获", [
        "技术收获:",
        "  ISTQB/ISO 29119-4 四种测试技术的实战应用",
        "  pydantic + streamlit + pytest 完整工具链",
        "  LLM prompt 工程: 结构化→风险→Oracle 三阶段中文 prompt",
        "工程收获:",
        "  core/ 与 UI/LLM 解耦的架构设计",
        "  TDD-first 开发 (先写测试再写实现)",
        "  交互式编辑的 Human-in-the-loop 设计模式",
        "  API 成本仅 USD 0.50, 工具辅助节省 64% 工时",
        "核心认知: AI 工具是测试工程师的倍增器, 而非替代品",
    ])

    # ── Slide 14: 改进建议 ──
    _add_content_slide(prs, "改进建议", [
        "工具改进:",
        "  可执行代码生成: 从 TestCase 自动生成 pytest 骨架",
        "  覆盖率反馈循环: pytest 结果回流到工具",
        "  多语言支持: 日/韩等语言的 prompt 模板",
        "AuthSystem 改进:",
        "  密码 bcrypt 哈希存储 (解决 R-09)",
        "  登录速率限制 (解决 R-01)",
        "  数据持久化: dict → SQLite",
        "  邮箱验证码机制 (增强注册安全)",
        "流程改进:",
        "  CI/CD 集成: GitHub Actions 自动运行流水线",
    ])

    # ── Slide 15: Q&A ──
    _add_title_slide(prs, "Q & A", "感谢观看 · Thank You\n\ngithub.com/your-org/autotestdesign")

    prs.save("demo/presentation.pptx")
    print("✓ demo/presentation.pptx 已生成")


if __name__ == "__main__":
    main()
